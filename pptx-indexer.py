from pptx import Presentation
import os
from pathlib import Path
import argparse
import re
import json
import subprocess
import sys
import logging

import PpIndexConfig as pic
from PpIndexCommon import remove_slides, replace_CSLandCSP

# デフォルト値
#default_source_trailer='_source'    # ソースファイルから除去する文字列
default_target_trailer='_generated'    # 生成されるファイル名に付加する文字列

# 使用する可能性のあるタグ
available_tagsregex = ['#CHAPT#', '#DT#']

# EaxyIndexer 用のデフォルト値
title_content_separator_EI=') '
number_delimiter_EI='.'

pptx_indexer_version='1.2.1'               # バージョン番号   
indexer_logformat='%(message)s'     # simple format
# ロガーの作成
indexer_logger = logging.getLogger(__name__)

# スクリプトのディレクトリ
script_dir = Path(__file__).resolve().parent
print(f"script_dir: {script_dir}")

easy_indexer_scriptname = 'PpEasyIndexer.py'
easy_indexer_scriptpath = script_dir / easy_indexer_scriptname
collector_scriptname = 'PpIndexCollector.py'
collector_scriptpath = script_dir / collector_scriptname 
labeler_scriptname = 'PpIndexLabeler.py'
labeler_scriptpath = script_dir / labeler_scriptname


# 初期化エラー例外を定義する
class InitializingError(Exception):
    def __init__(self, message):
        self.message = message
        super().__init__(message)

    def __str__(self):
        return f"InitializingError: {self.message}"


def setLogger(loglevel, logpath, logoutput):
    # ログレベルの設定
    if loglevel.upper() == 'DEBUG':
        indexer_logger.setLevel(logging.DEBUG)
    else:
        indexer_logger.setLevel(logging.INFO)

    # logoutput が STDOUT の場合は標準出力に、それ以外の場合はファイルに出力
    if logoutput.upper() == 'STDOUT':
        console_handler = logging.StreamHandler(stream=sys.stdout)
        console_handler.setFormatter(logging.Formatter(indexer_logformat))
        indexer_logger.addHandler(console_handler)
    else:
        file_handler = logging.FileHandler(logpath / logoutput, 'a')
        file_handler.setFormatter(logging.Formatter(indexer_logformat))
        indexer_logger.addHandler(file_handler)


# コマンドライン引数を解析して返す
def parse_commandargs():
    # コマンドライン引数のパーサーを作成
    parser = argparse.ArgumentParser(description="Generates new PowerPoint files with hierarchical indexed labels specified by DT or CHAPT tags.")
    # 'file' 引数の定義
    parser.add_argument("file", help="PPTX source file or YAML configuration file. Extensions can be omitted.")
    # 'file' 引数の定義
    parser.add_argument('--parameter', '-pa', action='store_true', help="indicates the file as YAML configuration file.")
    # ログレベル指定　オプション --ll 、略称 -l,文字型、デフォルトは info
    parser.add_argument('--loglevel', '-l', type=str, default='info', help='log level [debug|INFO]')
    # ログファイル指定　オプション --lf, 略称 -f,文字型、デフォルトは STDOUT
    parser.add_argument('--logfile', '-f', type=str, default='STDOUT', help='path to log file, or STDOUT if omitted')
    # オプション引数 --dump 、略称 -d,文字型、デフォルトは 空文字列 を指定
    parser.add_argument('--dump', '-d', type=str, default='', help='dump file name')
    # バージョン番号を表示
    parser.add_argument('--version', '-v', action='version', version=f'%(prog)s {pptx_indexer_version}')
    # 生成されるファイル (target) のファイル名に付加される文字列
    parser.add_argument('--target-trailer', '-tt', type=str, help='string appended to the generated file name')

    # ソースファイル (source) のファイル名に付加される文字列
    parser.add_argument('--remove-source-trailer',  '-rst', type=str,  help='remove source trailer from the generated file name')
    # 中間生成ファイルを削除するか否か
    parser.add_argument('--remove-workfile', '-r', action='store_true', help=f'remove working .json files after completion', )

    # ------- parameters for EasyIndexer, 簡易版向けのパラメータ ---------
    parser.add_argument('--skipsildes', '-ss', type=int, default=0, help='number of slides to be skipped at the beginning')
    parser.add_argument('--separator', '-sp', type=str, default=title_content_separator_EI, help='separator between chapter.section and title content')
    parser.add_argument('--delimitter', '-dl', type=str, default=number_delimiter_EI, help='delimitter between chapter and section numbers')
    parser.add_argument('--deletecsl', '-csl', action='store_true', help='delete slides containing #CSL#')
    parser.add_argument('--deletecsp', '-csp', action='store_true', help='delete shapes containing #CSP#')

    # 引数を解析
    _args = parser.parse_args()
    return _args



def find_Tags_in_pptx(file_path, tagsregex):
    # プレゼンテーションをロード
    prs = Presentation(file_path)

    # すべてのスライドをループ処理して連結テキストを作る
    combinedtext=''
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        combinedtext += run.text

    # テキストからタグを検索
    found = {}
    for tag in tagsregex:
        if re.search(tag, combinedtext):
            found[tag] = True

    return found


def onefile_collect_and_label(file_path,args):
    """ 単一のソースファイルに対してcollectorとlabelerを実行する
    """

    print("onefile_collect_and_label:", file_path)

    #----------------------------------------------------
    # collector を実行
    #----------------------------------------------------
    parameters = ['python', collector_scriptpath, file_path]
    parameters.append('-pp')    # パラメータファイルを使わない

    if args.loglevel:
        parameters.append('-l')
        parameters.append(args.loglevel)
    if args.logfile:
        parameters.append('-f')
        parameters.append(args.logfile)

    print("processing collector:", parameters)
    result = subprocess.run(parameters, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding='cp932')

    print(result.stdout)

    #----------------------------------------------------
    # labeler を実行
    #----------------------------------------------------
    parameters = ['python', labeler_scriptpath, file_path]
    parameters.append('-pp')    # パラメータファイルを使わない
    if args.loglevel:
        parameters.append('-l')
        parameters.append(args.loglevel)
    if args.logfile:
        parameters.append('-f')
        parameters.append(args.logfile)
    if args.deletecsl:
        parameters.append('-csl')
    if args.deletecsp:
        parameters.append('-csp')   
    if args.target_trailer:
        parameters.append('-tt')
        parameters.append(args.target_trailer)
    if args.remove_source_trailer:
        parameters.append('-rst')
        parameters.append(args.remove_source_trailer)

    # if args.remove_workfile:
    #     後でワークファイル削除の処理をここに組み込む

    print("processing labeler:", parameters)

    result = subprocess.run(parameters, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding='cp932')

    print(result.stdout)

    pass



def onefile_easyindexer(file_path,args):
    """ 単一のソースファイルに対してEasyIndexerを実行する

    Args:
        file_path (str): ソースファイルのパス(拡張子 .pptx つき)
        args (argparse.Namespace): コマンドライン解析結果

    Returns:
        なし
    
    """
    parameters = ['python', easy_indexer_scriptpath, file_path]
    if args.skipsildes:
        parameters.append('-ss')
        parameters.append(str(args.skipsildes))
    if args.separator:
        parameters.append('-sp')
        parameters.append(args.separator)
    if args.delimitter:
        parameters.append('-dl')
        parameters.append(args.delimitter)
    if args.deletecsl:
        parameters.append('-csl')
    if args.deletecsp:
        parameters.append('-csp')   
    if args.target_trailer:
        parameters.append('-tt')
        parameters.append(args.target_trailer)
    if args.remove_source_trailer:
        parameters.append('-rst')
        parameters.append(args.remove_source_trailer)

    # EasyIndexer don't use workfile
    # if args.remove_workfile:
    #     parameters.append('-r')

    if args.loglevel:
        parameters.append('-l')
        parameters.append(args.loglevel)
    if args.logfile:
        parameters.append('-f')
        parameters.append(args.logfile)

    print("onefile_easyindexer:", parameters)

    result = subprocess.run(parameters, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding='cp932')

    print(result.stdout)


#----------------------------------------------------
# yaml パラメータファイルを読み込んで collector, labeler により処理する
#----------------------------------------------------
def read_parameter_and_process(file_path,args):
    """ yaml パラメータファイルを読み込んで collector, labeler により処理する

    Args:
        file_path (str): パラメータファイルのパス(拡張子 .yaml つき)
        args (argparse.Namespace):  コマンドライン解析結果

    Returns:
        なし
    
    """

    #----------------------------------------------------
    # collector を実行
    #----------------------------------------------------
    parameters = ['python', collector_scriptpath, file_path]
    if args.loglevel:
        parameters.append('-l')
        parameters.append(args.loglevel)
    if args.logfile:
        parameters.append('-f')
        parameters.append(args.logfile)

    print("processing collector:", parameters)
    result = subprocess.run(parameters, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding='cp932')

    print(result.stdout)


    #----------------------------------------------------
    # labeler を実行
    #----------------------------------------------------
    parameters = ['python', labeler_scriptpath, file_path]
    if args.loglevel:
        parameters.append('-l')
        parameters.append(args.loglevel)
    if args.logfile:
        parameters.append('-f')
        parameters.append(args.logfile)

    print("processing labeler:", parameters)

    result = subprocess.run(parameters, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding='cp932')

    print(result.stdout)

    # if args.remove_workfile:
    #     後でワークファイル削除の処理をここに組み込む




#----------------------------------------------------
# 指定した extention を持つ、実在するファイルのパスを返す
#----------------------------------------------------
def getExistingfilePathExt( filepath , extension ):
    """ 指定した extention を持つ、実在するファイルのパスを返す

    Args:
        filename (str): ファイル名、拡張子がない場合は extension が付加される
        extension (str): 拡張子

    Returns:
        str: 実在するファイルのパス、拡張子つき。実在しない場合は None
    
    Examples:
        >>> getExistingfilePathExt('sample', '.pptx')
        'sample.pptx'
        >>> getExistingfilePathExt('sample.pptx', '.pptx')
        'sample.pptx'
        >>> getExistingfilePathExt('notexistssample.pptx', '.yaml')
        None

    """

    #filepathの拡張子を取得
    _ext = os.path.splitext(filepath)[1]

    if not _ext:    # 拡張子がない場合は extension を付加
        filepath += extension
    elif _ext != extension:   # 拡張子が異なる場合は None を返す
        return None

    if os.path.exists(filepath):   # 実在性確認
        return filepath
    else:
        return None



#----------------------------------------------------
## メイン処理 ##
if __name__ == "__main__":

    try:
        args = parse_commandargs()
        # print(args)
        # print(args.file)
        # print(args.parameter)
        # print(args.loglevel)
        # print(args.logfile)
        # print(args.dump)        
        # print(args.target_trailer)
        # print(args.source_trailer)
        # print(args.remove_workfile)
        # print(args.skipsildes)
        # print(args.separator)
        # print(args.delimitter)
        # print(args.deletecsl)
        # print(args.deletecsp)   

        if args.parameter:
            # file をパラメータファイルとして読み込む
            filepath = getExistingfilePathExt(args.file, '.yaml')
            if filepath is None:
                raise InitializingError(f"Incorrect parameter file specified: {args.file}")
            else:
                # パラメータファイルを読み込む。この場合は Collector, Labeler により処理する
                print(f"Parameter file: {filepath}")
                read_parameter_and_process(filepath,args)
        else:
            # file をソースファイルとして処理
            filepath = getExistingfilePathExt(args.file, '.pptx')
            if filepath is None:
                raise InitializingError(f"Incorrect source file specified: {args.file}")

            print(f"Source file: {filepath}")
            # ソースの実在性は確認されたので、タグを検索
            found = find_Tags_in_pptx(filepath, available_tagsregex)
            if '#CHAPT#' in found and '#DT#' in found:
                raise InitializingError("Can't use both #CHAPT# and #DT#")
            if not '#CHAPT#' in found and not '#DT#' in found:
                raise InitializingError("Use either #CHAPT# or #DT#")
            if '#CHAPT#' in found:
                # EasyIndexer により処理する
                print("found #CHAPT#")
                onefile_easyindexer(filepath,args)
            else:
                # Collector, Labeler により処理する
                print("found #DT#")
                onefile_collect_and_label(filepath,args)
    
    except Exception as e:
        print(e)
        exit()


