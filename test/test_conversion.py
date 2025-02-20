# PPTXコンバージョンのテスト
# collector と labeler をかけてpptxの SOURCE と TARGET を比較する

import pytest
import yaml
import logging
import subprocess
from pptx import Presentation
import re

logger = logging.getLogger('testlogger')
logger.setLevel(logging.DEBUG)
filehandler = logging.FileHandler('test/test.log')
logger.addHandler(filehandler)

# 比較対象とするテキストを抽出するための正規表現
# '!!S_99!!', '!! S_1 !!', '!! S_ A123 !!' などを抽出する。
# 空白があってもよく、空白部分を除去して比較のキーとして使う
extract_sequence = '!!( ?S_ ?[0-9a-zA-Z_]+) ?!!'

testdatafolder = 'test/testdata'


# yaml ファイルの単純な読込み。load_yamlで読み込んだデータとの照合用
def justreadyaml(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        _data = yaml.safe_load(file)
    return _data

def do_conversion(data):
    # data["CONFIG"] コンバージョンの設定ファイル
    # data["CONFIG"] を引数にして PpIndexCollector と PpIndexLabeler を実行
    #print(data["CONFIG"])
    result = subprocess.run(['python', 'PpIndexCollector.py', '--loglevel=info', '--logfile=collector_test.log', data["CONFIG"]], stdout=subprocess.PIPE)
    if result.returncode != 0:
        return False
    #print(f"result.stdout={result.stdout}")
    result = subprocess.run(['python', 'PpIndexLabeler.py', '--loglevel=info', '--logfile=collector_test.log', data["CONFIG"]], stdout=subprocess.PIPE)
    if result.returncode != 0:
        return False
    #print(f"result.stdout={result.stdout}")
    return True

# pptxファイルからtextを抽出し、paragraph 単位で連結して、!!\d+!! という文字列を含むものを配列として返す
def do_extraction(file_path):   
    #print(f"file_path={file_path}")
    prs = Presentation(file_path)

    ttdict = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    singleline = ''
                    for run in paragraph.runs:
                        singleline += run.text
                    # !!\d+!! という文字列を含むものを抽出する
                    cmpresult = re.search(extract_sequence, singleline)
                    if cmpresult:
                        #print(f"SL:{singleline}")
                        key = re.sub(r"\s+", "", cmpresult.group(1))
                        ttdict[key] = singleline
                        #print(f"■key={key}, singleline={singleline}")
    #print(f"■■ttdict={ttdict}")
    return ttdict    


def test_conversion():
    parameters = justreadyaml(f'{testdatafolder}/conversion_parameters.yaml')
    for i in range(len(parameters["DATA"])):
        #print(f"ID:{parameters['DATA'][i]['ID']}, PURPOSE:{parameters['DATA'][i]['PURPOSE']}")
        do_conversion(parameters["DATA"][i])
        # ここで、TARGET と VERIFIER を比較する
        #print(f"TARGET:{parameters['DATA'][i]['TARGET']}, VERIFIER:{parameters['DATA'][i]['VERIFIER']}")
        #print(f"extracting TARGET {parameters['DATA'][i]['TARGET']}")
        tglines = do_extraction(parameters['DATA'][i]['TARGET'])
        #print(f"extracting VERIFIER {parameters['DATA'][i]['VERIFIER']}")
        vflines = do_extraction(parameters['DATA'][i]['VERIFIER'])
        # print("-------------------------------")
        # for key in tglines.keys():
        #      print(f"  {key} : {tglines[key]}")
        # print("-------------------------------")
        # for key in vflines.keys():
        #      print(f"  {key} : {vflines[key]}")
        # print("-------------------------------")

        cmparefileresult = True
        for key in tglines.keys():
            match = (tglines[key] == vflines[key])
            if not match:
                cmparefileresult = False
                print(f"  differ T [{key}] {tglines[key]}")
                print(f"  differ V [{key}] {vflines[key]}")
            # else:
            #     print(f"  same [{key}] {tglines[key]}")
            #     print(f"  same [{key}] {vflines[key]}")

        assert cmparefileresult

