from pptx import Presentation
import win32com.client
import os
from pathlib import Path
import argparse
import re
import json
import sys
import logging

import PpIndexConfig as pic
from PpIndexCommon import remove_slides, replace_CSLandCSP, ProcessError,generate_source_and_target_filename

# 指定した .pptx ファイルとインデックスファイル(.json) を読んでタグを変換する

default_target_trailer='_generated'    # 生成されるファイル名に付加する文字列
default_source_trailer='_source'    # ソースファイルから除去する文字列

version='1.2'
logformat='%(message)s'                                     # simple format
# logformat='%(asctime)s - %(levelname)s - %(message)s'     # standard format


# ロガーの作成
logger = logging.getLogger(__name__)

def setLogger(loglevel, logpath, logoutput):
    # ログレベルの設定
    if loglevel.upper() == 'DEBUG':
        logger.setLevel(logging.DEBUG)
    else:
        logger.setLevel(logging.INFO)

    # logoutput が STDOUT の場合は標準出力に、それ以外の場合はファイルに出力
    if logoutput.upper() == 'STDOUT':
        console_handler = logging.StreamHandler(stream=sys.stdout)
        console_handler.setFormatter(logging.Formatter(logformat))
        logger.addHandler(console_handler)
    else:
        file_handler = logging.FileHandler(logpath / logoutput, 'a')
        file_handler.setFormatter(logging.Formatter(logformat))
        logger.addHandler(file_handler)



# ファイル/フォルダーが存在しない場合にエラーを発生させる
def raiseProcessError_if_not_exists(file_path):
    if not os.path.exists(file_path):
        raise ProcessError(f"File/Folder not found: {file_path}")


# コマンドライン引数を解析して返す
def parse_commandargs():
    # コマンドライン引数のパーサーを作成
    parser = argparse.ArgumentParser(description="Generates new PowerPoint files with labels converted according to the specified index files.")
    # ----------------------------------------------
    # yaml指定時とpptx単独指定時に共通する引数の定義
    parser.add_argument("file", help="PPTX source file or YAML configuration file. Extensions can be omitted.")
    # ログレベル指定　オプション --ll 、略称 -l,文字型、デフォルトは info
    parser.add_argument('--loglevel', '-l', type=str, default='info', help='log level [debug|INFO]')
    # ログファイル指定　オプション --lf, 略称 -f,文字型、デフォルトは STDOUT
    parser.add_argument('--logfile', '-f', type=str, default='STDOUT', help='path to log file, or STDOUT if omitted')
    # バージョン番号を表示
    parser.add_argument('--version', '-v', action='version', version=f'%(prog)s {version}')
    


    # ----------------------------------------------
    # pptx単独指定時用の定義
    parser.add_argument('--powerpoint', '-pp', action='store_true', help='indicates the file as PowerPoint source file.')
    parser.add_argument('--target-trailer', '-tt', type=str, help='string appended to the generated file name')
    parser.add_argument('--deletecsl', '-csl', action='store_true', help='delete slides containing #CSL#')
    parser.add_argument('--deletecsp', '-csp', action='store_true', help='delete shapes containing #CSP#')
    # ソースファイル (source) のファイル名に付加される文字列
    parser.add_argument('--remove-source-trailer',  '-rst', type=str,  help='remove source trailer from the generated file name')

    # 引数を解析
    _args = parser.parse_args()
    return _args


# indexファイルをもとに .pptx ファイルのラベル置換を行い、新しい .pptx ファイルを出力する
def generate_target(genparams,folderobj):
    _sourcepptx = genparams['SOURCE']
    _sourcepptxpath = folderobj / _sourcepptx
    _generatepptx = genparams['GENERATE']
    _generatepptxpath = folderobj / _generatepptx
    logger.info(f"source:{_sourcepptx} target:{_generatepptx}")

    prs = Presentation(_sourcepptxpath)

    # 削除対象スライドを削除して、スライド番号マッピング用ハッシュを受け取る key:元のスライド番号 value:新しいスライド番号
    snummap = remove_slides(prs, genparams['CSL'], logger)  

    firstlevelreplacements = []
    secondlevelreplacements = []
    mokujireplacements = []
    # 置換パターンを作成
    for indexfile in genparams['INDEX']:
        _indexpath = folderobj / indexfile
        logger.debug(f"indexfile:{indexfile} ")
        # jsonファイルを読み込む
        with open(_indexpath, 'r', encoding='utf-8') as file:
            indexdata = json.load(file)

        # 置換パターン
        # #DT# タイトル部のインデックス
        # #RI# 参照部　インデックス
        # #RT# 参照部　タイトル
        # #RN# 参照部　スライド番号
        # #RIT# 参照部　インデックスとタイトル
        # firstlevel の置換パターンを作成
        for key in indexdata:
            logger.debug(f"key:{key}")
            logger.debug(f"index:{indexdata[key]['index']}")
            #logger.debug(f"slidenumber:{snummap[indexdata[key]['slidenumber']]}")   # スライド削除後の番号へのマッピング対応
            logger.debug(f"slidenumber:{indexdata[key]['slidenumber']}")            # スライド削除後の番号へのマッピング対応はいったんやめる
            logger.debug(f"title:{indexdata[key]['title']}")
            # #DT#FA1# を indexdata[key]['index'] で置換するためのデータを作成
            replacements = {
                make_replacetargetSingleKey_reg("DT",key): rf"{indexdata[key]['index']}\g<1>",
                make_replacetargetSingleKey_reg("RI",key): rf"{indexdata[key]['index']}\g<1>",
                make_replacetargetSingleKey_reg("RT",key): rf"{indexdata[key]['title']}\g<1>",
                make_replacetargetSingleKey_all_reg("RST",key): rf"{indexdata[key]['title']}\g<1>",
                #make_replacetargetSingleKey_reg("RN",key): rf"{snummap[indexdata[key]['slidenumber']]}\g<1>",   # スライド削除後の番号へのマッピング対応
                make_replacetargetSingleKey_reg("RN",key): rf"{indexdata[key]['slidenumber']}\g<1>",                # スライド削除後の番号へのマッピング対応はいったんやめる
                make_replacetargetSingleKey_reg("RIT",key): rf"{indexdata[key]['index']} {indexdata[key]['title']}\g<1>",
                make_replacetargetSingleKey_all_reg("RSIT",key): rf"{indexdata[key]['index']} {indexdata[key]['title']}\g<1>",
            }
            # print(replacements)
            firstlevelreplacements.append(replacements)

            # secondlevel の置換パターンを作成
            secondlevelassoc = indexdata[key]['secondlevelassoc']
            secondlevelidxtitlelist = ""
            secondlevelslidenumberlist = ""
            logger.debug(f"secondlevelassoc:{secondlevelassoc}")
            logger.debug(f"snummap:{snummap}")
            for skey in secondlevelassoc:
                logger.debug(f"skey:{skey}")
                logger.debug(f"index:{secondlevelassoc[skey]['index']}")
                #logger.debug(f"slidenumber:{snummap[secondlevelassoc[skey]['slidenumber']]}")
                logger.debug(f"slidenumber:{secondlevelassoc[skey]['slidenumber']}")        # スライド削除後の番号へのマッピング対応はいったんやめる
                logger.debug(f"title:{secondlevelassoc[skey]['title']}")
                # #T#FA1.S1# を indexdata[key]['index'] で置換するためのデータを作成
                replacements = {
                    make_replacetargetDualKey_reg("DT",key, skey): rf"{indexdata[key]['index']}\g<1>{secondlevelassoc[skey]['index']}\g<2>",
                    make_replacetargetDualKey_reg("RI",key, skey): rf"{indexdata[key]['index']}\g<1>{secondlevelassoc[skey]['index']}\g<2>",
                    make_replacetargetDualKey_reg("RT",key, skey): rf"{secondlevelassoc[skey]['title']}\g<2>",
                    make_replacetargetDualKey_all_reg("RST",key, skey): rf"{secondlevelassoc[skey]['title']}\g<2>",
                    #make_replacetargetDualKey_reg("RN",key, skey): rf"{snummap[secondlevelassoc[skey]['slidenumber']]}\g<2>",
                    make_replacetargetDualKey_reg("RN",key, skey): rf"{secondlevelassoc[skey]['slidenumber']}\g<2>",        # スライド削除後の番号へのマッピング対応はいったんやめる
                    make_replacetargetDualKey_reg("RIT",key, skey): rf"{indexdata[key]['index']}\g<1>{secondlevelassoc[skey]['index']}\g<2> {secondlevelassoc[skey]['title']} ",
                    make_replacetargetDualKey_all_reg("RSIT",key, skey): rf"{indexdata[key]['index']}\g<1>{secondlevelassoc[skey]['index']}\g<2> {secondlevelassoc[skey]['title']} ",
                }
                # print(replacements)
                secondlevelreplacements.append(replacements)
                secondlevelidxtitlelist += f"{indexdata[key]['index']}.{secondlevelassoc[skey]['index']}) {secondlevelassoc[skey]['title']}\n"
                #secondlevelslidenumberlist += f"{snummap[secondlevelassoc[skey]['slidenumber']]}\n"
                secondlevelslidenumberlist += f"{secondlevelassoc[skey]['slidenumber']}\n"      # スライド削除後の番号へのマッピング対応はいったんやめる

            replacements = {
                make_replacetargetSingleKey_reg("RLISTIT",key): rf"{secondlevelidxtitlelist}\g<1>",
                make_replacetargetSingleKey_reg("RLISTSN",key): rf"{secondlevelslidenumberlist}\g<1>",
            }
            mokujireplacements.append(replacements)

    logger.debug('-----------------')
    logger.debug(firstlevelreplacements)
    logger.debug('-----------------')
    logger.debug(secondlevelreplacements)
    logger.debug('-----------------')
    logger.debug(mokujireplacements)

    logger.debug(f"fistlevelreplacements:{firstlevelreplacements}")

    # slidenumber はスライド削除後のスライド番号を保持
    slidenumber=0

    needshapedeletion = False

    for slide in prs.slides:
        logger.debug(f"slidenumber:{slidenumber}")
        slidenumber += 1
        tfnumber = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                tfnumber += 1
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        tobreak, _needshapedeletetion = replace_CSLandCSP(run, genparams['CSL'], genparams['CSP'],logger)
                        if _needshapedeletetion:
                            needshapedeletion = True
                        if tobreak:
                            break

                    cutoff = False
                    for run in paragraph.runs:
                        # run.text の置換を行う
                        if cutoff :
                            # #RST# または #RSIT# があったらそれに続くテキストは削除
                            run.text = ""
                        else:
                            (replacedtext, cutoff) = replace_text(run.text, firstlevelreplacements, secondlevelreplacements, mokujireplacements)
                            run.text = replacedtext


    # いったん変更を保存する
    prs.save(_generatepptxpath)

    # 絶対パスでないと win32 アプリでファイルを開けなかったので絶対パス取得
    newpptx_path_abs = os.path.abspath(_generatepptxpath)
    logger.debug(f"BEFORE win32 _generatepptxpath:{_generatepptxpath}")
    logger.debug(f"BEFORE win32 newpptx_path_abs :{newpptx_path_abs}")

    # シェイプ削除がある場合は pywin32 で削除処理を行う
    # (python-pptx でXMLエレメントをremoveする方法で削除しようとすると、保存後のpptxを開いたときに
    # 構成エラーが発生するので、pywin32 を使って削除します。ただしこの方法は遅い)
    if needshapedeletion:
        win32pptapp = win32com.client.Dispatch("Powerpoint.Application")
        win32prs = win32pptapp.Presentations.Open(str(newpptx_path_abs), WithWindow=False)
        for slide in win32prs.Slides:
            # 逆順でシェイプをループ（削除中にコレクションを変更しないように）
            for shape in reversed(list(slide.Shapes)):
                if shape.HasTextFrame == -1 and shape.TextFrame.HasText:
                    # テキストの中に #CSP# または #TEMP# または #CSL# の文字があれば削除
                    if "#CSP#" in shape.TextFrame.TextRange.Text or "#TEMP#" in shape.TextFrame.TextRange.Text or "#CSL#" in shape.TextFrame.TextRange.Text:
                        logger.debug(f"Found #CSP# in text box {shape.TextFrame.TextRange.Text}")
                        shape.Delete()
#                    pass

        logger.info(f"saving as newpptx_path_abs:{newpptx_path_abs}")
        win32prs.SaveAs(newpptx_path_abs)

        # プレゼンテーションを閉じる
        win32prs.Close()

        # PowerPointを閉じる
        # win32pptapp.Quit()

    # generate_target終了



def make_replacetargetDualKey_reg(tag,key, skey):
    # 「ピリオドまたはハイフンで区切ったキー文字列」の後ろに # と「非英数文字または空白または文末」が続く場合にマッチ
    return rf"#{tag}#{key}([-.]){skey}#([^0-9a-zA-Z_.]|\s|$)"

def make_replacetargetSingleKey_reg(tag,key):
    # 「キー文字列」の後ろに # と「非英数文字または空白または文末」が続く場合にマッチ
    return rf"#{tag}#{key}#([^0-9a-zA-Z_.]|\s|$)"

def make_replacetargetDualKey_all_reg(tag,key, skey):
    # 「ピリオドまたはハイフンで区切ったキー文字列」の後ろに # と「任意の文字列または文末」が続く場合にマッチ
    # "#RITS#CHAP.SECT# タイトル文字列" のようなタグを後のタイトル文字列も含めて置換する用途に使う
    return rf"#{tag}#{key}([-.]){skey}#.?(.*|$)"

def make_replacetargetSingleKey_all_reg(tag,key):
    # 「キー文字列」の後ろに # と「任意の文字列または文末」が続く場合にマッチ
    # "#RITS#CHAP# タイトル文字列" のようなタグを後のタイトル文字列も含めて置換する用途に使う
    return rf"#{tag}#{key}#.?(.*|$)"




# text に対して firstlevel と secondlevel の置換を行う
def replace_text(text, first, second,mokuji):
    # text に #\w+# が含まれない場合はそのまま返す
    cutoff = False
    if not re.search(r'#\w+#', text):
        return (text, cutoff)
    
    for replacements in second:
        for target, replacement in replacements.items():
            logger.debug(f"searching text:{text} target:{target} replacement:{replacement}")
            if re.search(target, text):
                if re.search(r'#RST#', text) or re.search(r'#RSIT#', text):
                    cutoff = True
                replacedtext = textreplacewrapper_reg(text, target, replacement)
                logger.debug(f"replacing text:{text} target:{target} replacement:{replacement}")
                logger.debug(f"  replacedtext:{replacedtext}")
                text = replacedtext

    for replacements in first:
        for target, replacement in replacements.items():
            logger.debug(f"searching text:{text} target:{target} replacement:{replacement}")
            if re.search(target, text):
                if re.search(r'#RST#', text) or re.search(r'#RSIT#', text):
                    cutoff = True
                replacedtext = textreplacewrapper_reg(text, target, replacement)
                logger.debug(f"replacing text:{text} target:{target} replacement:{replacement}")
                logger.debug(f"  replacedtext:{replacedtext}")
                text = replacedtext

    for replacements in mokuji:
        for target, replacement in replacements.items():
            #logger.debug(f"searching text:{text} target:{target} replacement:{replacement}")
            if re.search(target, text):
                if re.search(r'#RST#', text) or re.search(r'#RSIT#', text):
                    cutoff = True
                replacedtext = textreplacewrapper_reg(text, target, replacement)
                logger.debug(f"replacing text:{text} target:{target} replacement:{replacement}")
                logger.debug(f"  replacedtext:{replacedtext}")
                text = replacedtext

    return (text,cutoff)


# text に対して target を replacement に置換する
def textreplacewrapper(text, target, replacement):
    logger.debug(f"text:{text} target:{target} replacement:{replacement}")
    return text.replace(target, replacement)


# text に対して target を replacement に置換する、正規表現バージョン
def textreplacewrapper_reg(text, target, replacement):
    logger.debug(f"text:{text} target:{target} replacement:{replacement}")
    replaced_text = re.sub(target, replacement, text)
    logger.debug(f"replaced_text:{replaced_text}")
    return replaced_text



def generating_process(generatingarray, folderobj):
    for i in range(len(generatingarray)):
        # _generating[i] のパス名に拡張子を付与、ファイルの存在確認
        logger.debug(f"generating: {generatingarray[i]}")
        _generate = pic.add_extension(pic.remove_extension(generatingarray[i]['GENERATE']), 'pptx')
        generatingarray[i]['GENERATE'] = _generate
        _generatepathobj = folderobj / _generate
        if os.path.exists(_generatepathobj):
            # _generatepathobj のファイルが既に存在したら警告
            logger.info(f"Warning: overwriting existing pptx file: {_generatepathobj}")

        _source = pic.add_extension(pic.remove_extension(generatingarray[i]['SOURCE']), 'pptx')
        generatingarray[i]['SOURCE'] = _source
        _sourcepathobj = folderobj / _source
        # _sourcepath のファイルが存在しなかったらエラー
        raiseProcessError_if_not_exists(_sourcepathobj)

        for j in range(len(generatingarray[i]['INDEX'])):
            _index = pic.add_extension(pic.remove_extension(generatingarray[i]['INDEX'][j]), 'json')
            generatingarray[i]['INDEX'][j] = _index
            _indexpathobj = folderobj / _index
            # _indexpath のファイルが存在しなかったらエラー
            raiseProcessError_if_not_exists(_indexpathobj)

        # ラベリング処理をする
        logger.debug(f"generating:\n  sourcepath:{_sourcepathobj}\n  generatepath:{_generatepathobj}")
        generate_target(generatingarray[i],folderobj)


# メイン処理
def main():
    global default_target_trailer, default_source_trailer

    args = parse_commandargs()

    if args.powerpoint:

        # remove-source-trailer が指定されていて、かつ target-trailer が指定されていない場合は
        # target-trailer を空文字列にする
        if args.remove_source_trailer:
            default_source_trailer=args.remove_source_trailer
            if not args.target_trailer:
                default_target_trailer=''

        # 単独指定された pptx のフォルダーを取得
        _folder = Path(args.file).parent
        _folderobj = Path(_folder)

        # sourcefile_path と targetfile_path を確定する
        sourcefile_path, targetfile_path = generate_source_and_target_filename(args, default_source_trailer, default_target_trailer,logger)

        sourcefilename = Path(sourcefile_path).name
        targetfilename = Path(targetfile_path).name

        # ロガーの設定
        setLogger(args.loglevel, _folderobj ,  args.logfile)

        _generatingarray=[
            {'GENERATE':pic.add_extension(pic.remove_extension(targetfilename), 'pptx'),
             'SOURCE':sourcefilename, 
             'INDEX':[pic.add_extension(pic.remove_extension(sourcefilename), 'json')], 
             'CSL': args.deletecsl ,
             'CSP': args.deletecsp },
             ]
        generating_process(_generatingarray, _folderobj)
        logger.info(f"finished labeling of {args.file}")

    else:
        # YAMLファイルを読み込む
        configs = pic.verify_parameter_formats(pic.load_yaml(args.file,logger=logger),logger=logger)
        raiseProcessError_if_not_exists(configs['FOLDER'])

        _folder = configs['FOLDER']
        _folderobj = Path(_folder)

        # ロガーの設定
        setLogger(args.loglevel, _folderobj ,  args.logfile)

        logger.info( "--------- Labeler configuration --------")
        logger.info(f"configs: {configs}")

        logger.info( "----------- Labeler processing ---------")
        _generatingarray=configs['GENERATING']
        generating_process(_generatingarray, _folderobj)

        logger.info(f"finished labeling of {args.file}")



if __name__ == "__main__":
    try:
        main()
        exit(0)
    except pic.ConfigError as e:
        logger.info(f"ConfigError: {e}")
        print(f"ConfigError: {e}")
        exit(1)
    except ProcessError as e:
        logger.info(f"ProcessError: {e}")
        print(f"ProcessError: {e}")
        exit(2)

