from pptx import Presentation
import logging
import re
from pathlib import Path

# PowerPoint(.pptx) ファイル処理の共通関数



def remove_slides(prs, deletion,logger):
    logger.info(f"remove_slides:deletion:{deletion}")
    snummap = {}    

    _snums_to_remove = ()   # 削除するスライドID
    _newsnum=0                 # スライド番号(削除後)
    _sourcesnum=0           # スライド番号(削除前)
    for slide in prs.slides:
        _sdelflag=False     # このスライドは削除するよフラグ
        snummap[_sourcesnum] = _newsnum
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # run.text に #CSL# が含まれていたら、そのslideをまるごと削除する
                        if deletion :
                            if re.search(r'#CSL#', run.text):
                                # 削除するスライド番号を _sid_to_remove に追加
                                _snums_to_remove += (_sourcesnum,)
                                logger.debug(f"deleting slide:{_sourcesnum}")
                                _sdelflag=True
                                break
        if not _sdelflag:
            _newsnum += 1
        _sourcesnum += 1

    logger.debug(f"snummap:{snummap}")

    logger.debug(f"_snums_to_remove:{_snums_to_remove}")
    xml_slides = prs.slides._sldIdLst  # スライドのXMLリストへのアクセス
    logger.debug(f"xml_slides:{xml_slides}")
    # _sid_to_remove に含まれるスライドIDを逆順に削除する。まあ逆順じゃなくてもかまわんみたいだが
    for _sn in reversed(_snums_to_remove):
          # 削除したいスライドID
        logger.debug(f"removing slide:{xml_slides[_sn]}")
        xml_slides.remove(xml_slides[_sn])

    return snummap


def replace_CSLandCSP(run, deletecsl, deletecsp,logger):
    # run.text に #CSP# が含まれていたら shape を削除 または #CSP# の文字列のみを削除
    tobreak = False
    needshapedeletetion = False
    logger.debug(f'delecsl:{deletecsl},deletecsp:{deletecsp}')
    if  re.search(r'#CSP#', run.text):
        if deletecsp:
            needshapedeletetion = True
            logger.debug(f"need shape deletion:{run.text}")
            # このシェイプは後にまるごと削除されるので、
            # 呼び出し元に戻った後の replace_text が実行不要、
            # run の残りの処理も不要なので tobreak=True にしておく
            tobreak = True
        else:
            # 正規表現を使って '#CSP#\s?' を削除
            replacedtext = re.sub(r'#CSP#\s?', '', run.text)
            run.text = replacedtext
            #continue
    
    # run.text に #TEMP# が含まれていたら shape を削除
    if  re.search(r'#TEMP#', run.text) or re.search(r'#MEMO#', run.text):
        needshapedeletetion = True
        logger.debug(f"need shape deletion:{run.text}")
        # このシェイプは後にまるごと削除されるので、
        # 呼び出し元に戻った後の replace_text が実行不要、
        # run の残りの処理も不要なので tobreak=True にしておく
        tobreak = True
    
    # run.text に #CSL# が含まれていて、genparams['CSL'] がFalseの場合は #CSL# とそれに続く文字列のみを削除
    if  re.search(r'#CSL#', run.text):
        if deletecsl:
            # ページ削除は既に済んでいるのでここでは特にやることはない
            # というか実際にはここには来ないので下の tobreak=True は dead code だが、
            # 呼び出し元に戻った後の replace_text が実行不要、
            # run の残りの処理も不要なことを明示するために残しておく
            tobreak = True
        else:
            # ページ削除をしない場合は #CSL# とそれに続く文字列のみを削除
            # 正規表現を使って '#CSL#\s?.*' を削除
            replacedtext = re.sub(r'#CSL#\s?.*', '', run.text)
            run.text = replacedtext
            #continue

    return tobreak, needshapedeletetion



# 対象ファイル読み込みエラー、書き込みエラーなど
class ProcessError(Exception):
    def __init__(self, message):
        self.message = message
        super().__init__(message)

    def __str__(self):
        return f"ProcessError: {self.message}"


def generate_source_and_target_filename( args, default_source_trailer, default_target_trailer,logger):
    # 指定されたファイル名に .pptx がなければ追加する
    logger.debug(f'args.file:{args.file}')
    logger.debug(f'default_source_trailer:{default_source_trailer}')
    logger.debug(f'default_target_trailer:{default_target_trailer}')

    sourcefile_path = Path(args.file)
    if not sourcefile_path.suffix:
        sourcefile_path = sourcefile_path.with_suffix('.pptx')

    # sourcefile_path をもとに targetfile_path を生成するのでいったん targetfile_path に sourcefile_path をコピー
    targetfile_path = sourcefile_path

    # remove_source_trailer が指定されていて、
    # args.file の stem が default_source_trailer で終わっていれば除去する
    if args.remove_source_trailer:
        if targetfile_path.stem.endswith(default_source_trailer):
            targetfile_path = targetfile_path.with_stem(targetfile_path.stem[:-len(default_source_trailer)])

    # sourcefile_path に target_trailer を追加して targetfile_path にする
    target_trailer = args.target_trailer if args.target_trailer else default_target_trailer
    targetfile_path = targetfile_path.with_stem(f"{targetfile_path.stem}{target_trailer}")

    # targetfile_path と sourcefile_path が同じ場合は エラーとする
    if targetfile_path == sourcefile_path:
        # 実際は default_target_trailer ='' と書き換えない限りこれは起こらない
        raise ProcessError('ソースファイルとターゲットファイルが同じです')

    return sourcefile_path, targetfile_path

    # この段階で sourcefile_path と targetfile_path が確定している
    