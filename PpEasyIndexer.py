#!/usr/bin/env python   # PowerPoint Easy Numbering tool
# -*- coding: utf-8 -*- 

from pptx import Presentation
import os
from pathlib import Path
import re
import argparse
import win32com.client
import logging
import sys

from PpIndexCommon import remove_slides, replace_CSLandCSP, ProcessError,generate_source_and_target_filename

#指定したPowerPointファイルを読み込み、章番号と節番号を付与して新しいファイルを生成します。
#章番号は #CHAPT# 、節番号は #SECTION# という文字列をスライドのタイトルに含むことで判別します。

version='1.2'
logformat='%(message)s'                                     # simple format
# logformat='%(asctime)s - %(levelname)s - %(message)s'     # standard format

# ロガーの作成
logger = logging.getLogger(__name__)

chapter_marker=r'#CHAPT#'
section_marker=r'#SECTION#'
chapter_title_format=r'#CHAPTNUM##SEPA##CONTENT#'
section_title_format=r'#CHAPTNUM##DELM##SECTNUM##SEPA##CONTENT#'

title_content_separator=') '
number_delimiter='.'
should_delete_csl=False
should_delete_csp=False

default_source_trailer='_source'    # ソースファイルから除去する文字列
default_target_trailer='_generated'    # 生成されるファイル名に付加する文字列

# chaptnumber=0
# sectnumber=0



def setLogger(loglevel, logpath, logoutput):
    # ログレベルの設定
    if loglevel.upper() == 'DEBUG':
        logger.setLevel(logging.DEBUG)
    else:
        logger.setLevel(logging.INFO)

    # logoutput が STDOUT の場合は標準出力に、それ以外の場合はファイルに出力
    if logoutput.upper() == 'STDOUT':
        console_handler = logging.StreamHandler(stream=sys.stdout)
        console_handler.setFormatter(logging.Formatter(logformat))
        logger.addHandler(console_handler)
    else:
        file_handler = logging.FileHandler(logpath / logoutput, 'a')
        file_handler.setFormatter(logging.Formatter(logformat))
        logger.addHandler(file_handler)

# コマンドライン引数を解析して返す
def parse_commandargs():
    # コマンドライン引数のパーサーを作成
    parser = argparse.ArgumentParser(description="A simple chapter and section numbering processor.")
    # 'pptx' 引数の定義　処理対象のpptxファイルを指定する
    parser.add_argument("file", help="PowerPoint file to process")
    parser.add_argument('--skipsildes', '-ss', type=int, default=0, help='number of slides to be skipped at the beginning')
    parser.add_argument('--separator', '-sp', type=str, default=title_content_separator, help='separator between chapter.section and title content')
    parser.add_argument('--delimitter', '-dl', type=str, default=number_delimiter, help='delimitter between chapter and section numbers')
    parser.add_argument('--deletecsl', '-csl', action='store_true', help='delete slides containing #CSL#')
    parser.add_argument('--deletecsp', '-csp', action='store_true', help='delete shapes containing #CSP#')
    parser.add_argument('--target-trailer',  '-tt', type=str,  help='string appended to the generated file name')
    parser.add_argument('--remove-source-trailer',  '-rst', type=str,  help='remove source trailer from the generated file name')
    # ログレベル指定　オプション --ll 、略称 -l,文字型、デフォルトは info
    parser.add_argument('--loglevel', '-l', type=str, default='info', help='log level [debug|INFO]')
    # ログファイル指定　オプション --lf, 略称 -f,文字型、デフォルトは STDOUT
    parser.add_argument('--logfile', '-f', type=str, default='STDOUT', help='path to log file, or STDOUT if omitted')
    # バージョン番号を表示
    parser.add_argument('--version', '-v', action='version', version=f'%(prog)s {version}')
    

    # 引数を解析
    _args = parser.parse_args()
    return _args


def modify_slide_titles(args):
    global chapter_marker,section_marker,chapter_title_format,section_title_format
    global title_content_separator,number_delimiter
    global should_delete_csl,should_delete_csp
    global default_target_trailer,default_source_trailer

    # sourcefile_path と targetfile_path を確定する
    sourcefile_path, targetfile_path = generate_source_and_target_filename(args, default_source_trailer, default_target_trailer,logger)

    skipsildes = args.skipsildes

    chaptnumber=0  # 章番号をインクリメント
    sectnumber=0    # 節番号をリセット（次スライドから1になるのでここでは0にする）

    prs = Presentation(sourcefile_path)  # PowerPointファイルを読み込む

    # #CSL# ページ削除オプションが有効な場合はここで行う
    if should_delete_csl:
        remove_slides(prs, should_delete_csl, logger)  

    titles = []
    startchapter=False
    startsection=False
    insection=False
    slidenumber=0

    shapedeletioncounter =0

    # 全スライドについてループ        
    for slide in prs.slides:
        slidenumber+=1
        if slidenumber <= skipsildes:
            logger.debug(f"[{slidenumber}] スキップ")
            continue

        title = ''
        # タイトルプレースホルダーを探す
        titleshape = None

        if hasattr(slide.shapes.title, 'has_text_frame'):
            logger.info(f"[{slidenumber}] 変換前タイトル:{slide.shapes.title.text}")
            title = slide.shapes.title.text
            titleshape = slide.shapes.title
        else:
            logger.info(f"[{slidenumber}] 変換処理なし")



        # 章節番号処理のための shapes ループ（#CSP#と#CSL#は別ループなことに注意）
        # スライド中に matchchapt または matchsect があるかどうかを確認
        startchapter=False
        startsection=False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    # matchchat または matchsect にマッチするtextを探す
                    if re.search(chapter_marker, text):
                        startchapter=True
                        text = text.replace(chapter_marker,'')  # タグを除去
                        run.text = text
                    if re.search(section_marker, text):
                        startsection=True
                        text = text.replace(section_marker,'')  # タグを除去
                        run.text = text
            # 両方が見つかったら終了。片方だけなら継続
            # if startchapter and startsection:
            #     break

        # ループ終了時には startchapter,startsection の判別がついている

        # タイトル変数はタグが除去されていないので除去する
        title = title.replace(chapter_marker,'') 
        title = title.replace(section_marker,'') 

        # 章・節番号の切り替わり処理.
        if startchapter and startsection:    
            # 両方発見されたので、章・節番号切り替わりの処理をする
            chaptnumber+=1  # 章番号をインクリメント
            sectnumber=1    # 節番号をリセット（このスライドから1にする）
            insection=True  # 次スライド以降はstartchapter,startsection未該当でも節が継続する
            newtitle = section_title_format
            newtitle = newtitle.replace('#CHAPTNUM#',str(chaptnumber))
            newtitle = newtitle.replace('#DELM#',number_delimiter)
            newtitle = newtitle.replace('#SECTNUM#',str(sectnumber))
            newtitle = newtitle.replace('#SEPA#',title_content_separator)
            newtitle = newtitle.replace('#CONTENT#',title)

        elif startchapter:    
            # matchchapt のみなので章番号切り替わりの処理のみ行う
            chaptnumber+=1  # 章番号をインクリメント
            sectnumber=0    # 節番号をリセット（次スライドから1になるのでここでは0にする）
            insection=True  # 次スライド以降はstartchapter,startsection未該当でも節が継続する
            newtitle = chapter_title_format
            newtitle = newtitle.replace('#CHAPTNUM#',str(chaptnumber))
            newtitle = newtitle.replace('#SEPA#',title_content_separator)
            newtitle = newtitle.replace('#CONTENT#',title)

        elif startsection or insection:    # 節番号切り替わりの処理のみ行う
            # matchsect のみなので節番号切り替わりの処理のみ行う
            if startsection:
                sectnumber=1    # 節番号をリセット（このスライドから1にする）
                insection=True  # 次スライド以降はstartchapter,startsection未該当でも節が継続する
            else:
                sectnumber+=1   # insection に該当するので節番号をインクリメント
            newtitle = section_title_format
            newtitle = newtitle.replace('#CHAPTNUM#',str(chaptnumber))
            newtitle = newtitle.replace('#DELM#',number_delimiter)
            newtitle = newtitle.replace('#SECTNUM#',str(sectnumber))
            newtitle = newtitle.replace('#SEPA#',title_content_separator)
            newtitle = newtitle.replace('#CONTENT#',title)

        else: 
            # 最初の Chapter より前のスライドはどこにも該当しない
            newtitle = title
            pass

        if titleshape:
            titleshape.text = newtitle
        titles.append(newtitle)

        logger.info(f"[{slidenumber}] 変換後タイトル:{newtitle}")


        # #CSP# と #CSL# の処理のための shapes ループ（章節番号処理とは別ループなことに注意）
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        tobreak, _needshapedeletetion = replace_CSLandCSP(run, should_delete_csl, should_delete_csp,logger)
                        logger.debug(f"tobreak:{tobreak},needshapedeletetion:{_needshapedeletetion}")
                        if _needshapedeletetion:
                            shapedeletioncounter += 1
                        if tobreak: # このシェイプの残りの run は処理不要
                            break

    # 新しいファイル名を生成
    # base_name, ext = os.path.splitext(sourcefile_path)
    # _generatepptxpath = f"{base_name}{target_trailer}{ext}"
    
    # 変更を保存
    prs.save(targetfile_path)


    # 絶対パスでないと win32 アプリでファイルを開けなかったので絶対パス取得
    newpptx_path_abs = os.path.abspath(targetfile_path)
    # print(f"BEFORE win32 _generatepptxpath:{_generatepptxpath}")
    # print(f"BEFORE win32 newpptx_path_abs :{newpptx_path_abs}")

    # シェイプ削除がある場合は pywin32 で削除処理を行う
    # (python-pptx でXMLエレメントをremoveする方法で削除しようとすると、保存後のpptxを開いたときに
    # 構成エラーが発生するので、pywin32 を使って削除します。ただしこの方法は遅い)
    if shapedeletioncounter > 0:
        logger.debug(f"Deleting shapes containing #CSP# using win32com")
        win32pptapp = win32com.client.Dispatch("Powerpoint.Application")
        win32prs = win32pptapp.Presentations.Open(str(newpptx_path_abs), WithWindow=False)
        for slide in win32prs.Slides:
            # 逆順でシェイプをループ（削除中にコレクションを変更しないように・・・まあ逆順じゃなくても大丈夫なようですが）
            for shape in reversed(list(slide.Shapes)):
                if shape.HasTextFrame == -1 and shape.TextFrame.HasText:
                    # テキストの中に #CSP# または #CSL# または #TEMP# または #MEMO# の文字があれば削除
                    text = shape.TextFrame.TextRange.Text
                    if "#CSP#" in text or "#CSL#" in text or "#TEMP#" in text or "#MEMO#" in text:
                        logger.debug(f"Deleting text box with CSP,CSL,TEMP,MEMO {shape.TextFrame.TextRange.Text}")
                        shape.Delete()

        win32prs.SaveAs(newpptx_path_abs)

        # プレゼンテーションを閉じる
        win32prs.Close()

        # PowerPointを閉じる
        #win32pptapp.Quit()


    logger.info(f"変換結果: {newpptx_path_abs}")


if __name__ == '__main__':

    try:    
        args = parse_commandargs()
        title_content_separator=args.separator
        number_delimiter=args.delimitter
        should_delete_csl=args.deletecsl
        should_delete_csp=args.deletecsp
        # remove-source-trailer が指定されていて、かつ target-trailer が指定されていない場合は
        # target-trailer を空文字列にする
        if args.remove_source_trailer:
            default_source_trailer=args.remove_source_trailer
            if not args.target_trailer:
                default_target_trailer=''

        _folder = Path(args.file).parent
        _folderobj = Path(_folder)

        # ロガーの設定
        setLogger(args.loglevel, _folderobj ,  args.logfile)

        modify_slide_titles(args)

    except Exception as e:
        logger.error(f"Error: {e}")
        sys.exit(1)

