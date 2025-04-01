from pptx import Presentation
import os
from pathlib import Path
import argparse
import re
import json
import sys
import logging

import PpIndexConfig as pic
from PpIndexCommon import remove_slides, ProcessError

# 指定した .pptx ファイルからインデックス用コードを抽出し、index ファイルを出力する

default_target_trailer='_generated'    # 生成されるファイル名に付加する文字列

version='1.2'
logformat='%(message)s'                                     # simple format
# logformat='%(asctime)s - %(levelname)s - %(message)s'     # standard format

# ロガーの作成
logger = logging.getLogger(__name__)

def setLogger(loglevel, logpath, logoutput):
    # ログレベルの設定
    if loglevel.upper() == 'DEBUG':
        logger.setLevel(logging.DEBUG)
    else:
        logger.setLevel(logging.INFO)

    # logoutput が STDOUT の場合は標準出力に、それ以外の場合はファイルに出力
    if logoutput.upper() == 'STDOUT':
        console_handler = logging.StreamHandler(stream=sys.stdout)
        console_handler.setFormatter(logging.Formatter(logformat))
        logger.addHandler(console_handler)
    else:
        file_handler = logging.FileHandler(logpath / logoutput, 'a')
        file_handler.setFormatter(logging.Formatter(logformat))
        logger.addHandler(file_handler)


# 「#DT#FA1) タイトル」 にマッチして FA1 と タイトル を取得
#level1_pattern = re.compile(r'#DT#([0-9a-zA-Z_]+)([^0-9a-zA-Z_.\s]+|\s?)\s+(.+)')
level1_pattern = re.compile(r'#DT#([0-9a-zA-Z_]+)#([^0-9a-zA-Z_.\s]+|\s?)\s+(.+)')
# ↑ 英数字1文字以上、[英数字|アンダースコア|ピリオド]以外の文字0文字以上、空白1文字以上、任意の文字列

# 「#DT#FA1.SZ1) タイトル」 にマッチして FA1 と SZ1 と タイトル を取得
level2_pattern = re.compile(r'#DT#([0-9a-zA-Z_]+)[-.]+([0-9a-zA-Z_]+)#([^0-9a-zA-Z_.\s]+|\s?)\s+(.+)')
# ↑ 英数字1文字以上、[英数字|アンダースコア]1文字以上、[ハイフン|ピリオド]1文字以上、[英数字|アンダースコア]1文字以上、[英数字|アンダースコア|ピリオド|空白]以外の文字0文字以上、空白1文字以上、任意の文字列

# 「#SUM# コンテンツ」 にマッチして コンテンツを取得
summary_pattern = re.compile(r'#SUM#\s+(\S.*)')
# ↑ 1文字以上の空白の後に任意の文字列

# グローバル変数
firstlevelassoc = {}
firstlevelnumber = 0        # 1階層目の連想配列のインデックス
slidenumber = 0             # スライド番号をカウントする




# ファイル/フォルダーが存在しない場合にエラーを発生させる
def raiseProcessError_if_not_exists(file_path):
    if not os.path.exists(file_path):
        raise ProcessError(f"File/Folder not found: {file_path}")


# コマンドライン引数を解析して返す
def parse_commandargs():
    # コマンドライン引数のパーサーを作成
    parser = argparse.ArgumentParser(description="Collects index information from PowerPoint files and outputs them as .json files.")
    # ----------------------------------------------
    # yaml指定時とpptx単独指定時に共通する引数の定義
    parser.add_argument("file", help="PPTX source file or YAML configuration file. Extensions can be omitted.")
    # ログレベル指定　オプション --ll 、略称 -l,文字型、デフォルトは info
    parser.add_argument('--loglevel', '-l', type=str, default='info', help='log level [debug|INFO]')
    # ログファイル指定　オプション --lf, 略称 -f,文字型、デフォルトは STDOUT
    parser.add_argument('--logfile', '-f', type=str, default='STDOUT', help='path to log file, or STDOUT if omitted')
    # バージョン番号を表示
    parser.add_argument('--version', '-v', action='version', version=f'%(prog)s {version}')
    # オプション引数 --dump 、略称 -d,文字型、デフォルトは 空文字列 を指定
    parser.add_argument('--dump', '-d', type=str, default='', help='dump file name')

    # ----------------------------------------------
    # pptx単独指定時用の定義
    parser.add_argument('--powerpoint', '-pp', action='store_true', help='indicates the file as PowerPoint source file.')
    parser.add_argument('--target-trailer', '-tt', type=str,  default=default_target_trailer ,help='string appended to the generated file name')
    parser.add_argument('--deletecsl', '-csl', action='store_true', help='delete slides containing #CSL#')
    parser.add_argument('--deletecsp', '-csp', action='store_true', help='delete shapes containing #CSP#')


    # 引数を解析
    _args = parser.parse_args()
    return _args



# .pptx ファイルを解析してインデックス対象テキストを抽出する
def collectandsave_index(indexparams, folderobj):
    global level1_pattern, level2_pattern, summary_pattern
    global firstlevelassoc, firstlevelnumber, slidenumber

    logger.info(f"Collecting indexes\n  from:{indexparams['SOURCE']}\n    to:{indexparams['INDEX']}")

    _sourcepptx = indexparams['SOURCE']
    _sourcepptxpath = folderobj / _sourcepptx
    _indexjson = indexparams['INDEX']
    _indexjsonpath = folderobj / _indexjson

    # .pptxファイルの読み込み
    prs = Presentation(_sourcepptxpath)

    # 削除対象スライドを削除して、スライド番号マッピング用ハッシュを受け取る key:元のスライド番号 value:新しいスライド番号
    snummap = remove_slides(prs, indexparams['CSL'],logger)

    firstlevelassoc = {}
    firstlevelnumber = 0
    # secondlevelnumber = 0     # ループ内でfirstlevelassocの中で初期化するので、ここでは不要
    slidenumber = 0

    # スライドごとにテキストを置換
    for slide in prs.slides:
        # note と summary はスライドごとに初期化
        notes_text = ''
        summary_text = ''

        # note があれば取得
        if slide.has_notes_slide:
            notes = slide.notes_slide
            notes_text = notes.notes_text_frame.text

        # 先にサマリーを取得
        for shape in slide.shapes:
            if shape.has_text_frame:
                frametext=''    # テキストフレームの全パラグラフを連結する変数
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        frametext += run.text
                matchsum = summary_pattern.search(frametext)
                if matchsum:
                    logger.debug(f'■matchsum■:{matchsum}')    
                    summary_text += matchsum.group(1)
                    logger.debug(f'summary:{summary_text}')

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    ptext=''    # 1パラグラフに含まれるテキストを連結してから、正規表現でマッチングする
                    for run in paragraph.runs:
                        ptext += run.text
                    #logger.debug(ptext)
                    match1 = level1_pattern.findall(ptext)      # 'FA1) タイトル' にマッチ
                    match2 = level2_pattern.findall(ptext)      # 'FA1.SZ1) タイトル' にマッチ
                    if match1 or match2:
                        logger.debug(ptext)
                    if match1:
                        firstleveltext = match1[0][0]
                        separator = match1[0][1]
                        titletext = match1[0][2]
                        logger.debug(f'match1 first:{firstleveltext} title:{titletext}')
                        logger.debug(f'firstlevelassoc:{firstlevelassoc}')
                        if firstleveltext not in firstlevelassoc:
                            addfirstlevel(firstleveltext, separator, titletext, notes_text, summary_text)
                            secondlevelnumber =0

                    if match2:
                        firstleveltext = match2[0][0]
                        secondleveltext = match2[0][1]
                        separator = match2[0][2]
                        titletext = match2[0][3]

                        if firstleveltext not in firstlevelassoc:
                            #　1階層目の定義がないまま2階層目が出現した場合は、1階層目を追加して処理を続行
                            # エラーをassertする・・・のはやめて、firstlevel を登録して処理続行
                            #assert False, f'firstleveltext:{firstleveltext} is not found in firstlevelassoc'
                            addfirstlevel(firstleveltext, separator, titletext, notes_text, summary_text)
                            secondlevelnumber =0    # 1階層目が追加されたので、2階層目のインデックスをリセット

                        # 1階層目で持っている連想配列をいったん取得
                        secondlevelassoc = firstlevelassoc[firstleveltext]['secondlevelassoc']
                        logger.debug(f"secondlevelassoc:{secondlevelassoc}")
                        if secondleveltext not in secondlevelassoc:
                            # 2階層目としてまだ出現していない文字列なら、既存のsecondlevelnumber + 1 で登録
                            logger.debug("secondleveltext not found in secondlevelassoc")
                            # ↓既存の2階層目のインデックスのリストを取得
                            secondlevelnumber_values = [secondlevelassoc[secondleveltext]['index'] for secondleveltext in secondlevelassoc if 'index' in secondlevelassoc[secondleveltext]]
                            logger.debug('secondlevelnumber_values:' + ', '.join(str(num) for num in secondlevelnumber_values))

                            secondlevelassoc[secondleveltext] = {}
                            # secondlevelnumber += 1
                            secondlevelnumber =  max(secondlevelnumber_values) + 1 if secondlevelnumber_values else 1   # 既存の最大値 + 1 で登録
                            # ↓1階層目の値を記録
                            secondlevelassoc[secondleveltext]['index'] = secondlevelnumber
                            secondlevelassoc[secondleveltext]['slidenumber'] = slidenumber
                            secondlevelassoc[secondleveltext]['separator'] = separator
                            secondlevelassoc[secondleveltext]['title'] = titletext
                            secondlevelassoc[secondleveltext]['notes'] = notes_text
                            secondlevelassoc[secondleveltext]['summary'] = summary_text
                            firstlevelassoc[firstleveltext]['secondlevelassoc'] = secondlevelassoc  # 1階層目の連想配列に書き戻す

                        logger.debug(f'match2 first:{firstleveltext} second:{secondleveltext} title:{titletext}')

        logger.debug("-1-----------------------------------------------")
        slidenumber += 1

    logger.debug("-2-----------------------------------------------")

    logger.debug(firstlevelassoc)

    logger.debug(f"indexing:\n  sourcepath:{_sourcepptxpath}\n  indexpath:{_indexjsonpath}")
    with open(_indexjsonpath, 'w', encoding='utf-8') as f:
        json.dump(firstlevelassoc, f, ensure_ascii=False, indent=4)
        logger.debug(f"index file saved: {_indexjsonpath}")


def addfirstlevel(firstleveltext, separator, titletext, notes_text, summary_text):
    global firstlevelassoc, firstlevelnumber, slidenumber
    firstlevelnumber += 1
    firstlevelassoc[firstleveltext] = {}
    # ↓1階層目の値を記録
    firstlevelassoc[firstleveltext]['index'] = firstlevelnumber
    firstlevelassoc[firstleveltext]['slidenumber'] = slidenumber
    firstlevelassoc[firstleveltext]['separator'] = separator
    firstlevelassoc[firstleveltext]['title'] = titletext
    firstlevelassoc[firstleveltext]['notes'] = notes_text
    firstlevelassoc[firstleveltext]['summary'] = summary_text
    # ↓2階層目の連想配列を用意しておく
    firstlevelassoc[firstleveltext]['secondlevelassoc'] = {}
    logger.debug(f'■■firstlevelnumber:{firstlevelnumber} summary:{summary_text}')
    logger.debug(firstlevelassoc)


def indexing_process(indexingarray, folderobj):
    for i in range(len(indexingarray)):
        # _indexing[i] のパス名に拡張子を付与、ファイルの存在確認
        logger.debug(f"indexing: {indexingarray[i]}")
        _index = pic.add_extension(pic.remove_extension(indexingarray[i]['INDEX']), 'json')
        indexingarray[i]['INDEX'] = _index
        _indexpathobj = folderobj / _index
        if os.path.exists(_indexpathobj):
            # _indexpath のファイルが既に存在したら警告
            logger.info(f"Warning: overwriting existing index file: {_indexpathobj}")

        _source = pic.add_extension(pic.remove_extension(indexingarray[i]['SOURCE']), 'pptx')
        indexingarray[i]['SOURCE'] = _source
        _sourcepathobj = folderobj / _source
        # _sourcepath のファイルが存在しなかったらエラー
        raiseProcessError_if_not_exists(_sourcepathobj)

        # インデックス生成処理をする
        logger.debug(f"indexing:\n  sourcepath:{_sourcepathobj}\n  indexpath:{_indexpathobj}")
        try:
            collectandsave_index(indexingarray[i], folderobj)
        except Exception as e:
            logger.error(f"Error collecting index from {_sourcepathobj}: {e}")
            raise ProcessError(f"Error collecting index from {_sourcepathobj}: {e}")

        logger.debug("-3-----------------------------------------------")


# メイン処理
def main():
    args = parse_commandargs()

    if args.powerpoint:

        # 単独指定された pptx のフォルダーを取得
        _folder = Path(args.file).parent
        _folderobj = Path(_folder)

        # args.file からファイル名のみ取得する
        sourcefilename = Path(args.file).name

        # ロガーの設定
        setLogger(args.loglevel, _folderobj ,  args.logfile)

        _indexingarray=[
            {'SOURCE':sourcefilename, 
             'INDEX':pic.add_extension(pic.remove_extension(sourcefilename), 'json'), 
             'CSL': args.deletecsl },
             ]
        indexing_process(_indexingarray, _folderobj)
        logger.info(f"finished collecting index information of {args.file}")

    else:
        # YAMLファイルを読み込む
        configs = pic.verify_parameter_formats(pic.load_yaml(args.file, dump=args.dump, logger=logger), logger=logger)
        raiseProcessError_if_not_exists(configs['FOLDER'])

        _folder = configs['FOLDER']
        _folderobj = Path(_folder)

        # ロガーの設定
        setLogger(args.loglevel, _folderobj ,  args.logfile)

        logger.info( "--------- Collector configuration --------")
        logger.info(f"configs: {configs}")

        logger.info( "----------- Collector  processing ---------")
        _indexingarray=configs['INDEXING']   
        indexing_process(_indexingarray, _folderobj)

        logger.info(f"finished collecting index information of {args.file}")


#----------------------------------------------
# メイン処理
if __name__ == "__main__":
    try:
        main()
        exit(0)
    except pic.ConfigError as e:
        logger.info(f"ConfigError: {e}")
        print(f"ConfigError: {e}")
        exit(1)
    except ProcessError as e:
        logger.info(f"ProcessError: {e}")
        print(f"ProcessError: {e}")
        exit(2)
